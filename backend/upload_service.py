
import io
import fitz 
import docx
import pandas as pd
from pptx import Presentation

MAX_CONTEXT_CHARS = 8000  


def from_pdf(content):
    doc = fitz.open(stream=content, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)

def from_docx(content):
    d = docx.Document(io.BytesIO(content))
    return "\n".join(p.text for p in d.paragraphs)

def from_csv(content):
    df = pd.read_csv(io.BytesIO(content))
    return df.to_string(index=False)

def from_txt(content):
    return content.decode("utf-8", errors="ignore")

def from_pptx(content):
    prs = Presentation(io.BytesIO(content))
    lines = [
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]
    return "\n".join(lines)


EXTRACTORS = {
    "pdf": from_pdf,
    "docx": from_docx,
    "csv": from_csv,
    "txt": from_txt,
    "pptx": from_pptx,
}


def extract_text(filename, content):
    ext = filename.lower().split(".")[-1]
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return None
    text = extractor(content).strip()
    return text[:MAX_CONTEXT_CHARS]