import io
import os
import chromadb
import fitz  # pymupdf
import docx
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv
from langchain_openai import OpenAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

load_dotenv()

embeddings = OpenAIEmbeddings(
    model="nvidia/llama-nemotron-embed-vl-1b-v2:free",
    base_url="https://openrouter.ai/api/v1",
    api_key=os.getenv("OPENROUTER_API_KEY"),
    check_embedding_ctx_length=False,
    model_kwargs={"encoding_format": "float"},
)

chroma_client = chromadb.PersistentClient(path="./chroma_db")
upload_collection = chroma_client.get_or_create_collection(name="user_uploads")

splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)


def _from_pdf(content):
    doc = fitz.open(stream=content, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)

def _from_docx(content):
    d = docx.Document(io.BytesIO(content))
    return "\n".join(p.text for p in d.paragraphs)

def _from_csv(content):
    df = pd.read_csv(io.BytesIO(content))
    return df.to_string(index=False)

def _from_txt(content):
    return content.decode("utf-8", errors="ignore")

def _from_pptx(content):
    prs = Presentation(io.BytesIO(content))
    lines = [
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]
    return "\n".join(lines)


EXTRACTORS = {
    "pdf": _from_pdf,
    "docx": _from_docx,
    "csv": _from_csv,
    "txt": _from_txt,
    "pptx": _from_pptx,
}


def extract_text(filename, content):
    ext = filename.lower().split(".")[-1]
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return None
    return extractor(content).strip()


def store_upload(cid, filename, content):
    text = extract_text(filename, content)
    if not text:
        return 0

    chunks = splitter.split_text(text)
    ids = [f"{cid}_{filename}_{i}" for i in range(len(chunks))]
    metadatas = [{"Cid": cid, "source": filename} for _ in chunks]
    vectors = embeddings.embed_documents(chunks)

    upload_collection.add(
        ids=ids,
        documents=chunks,
        metadatas=metadatas,
        embeddings=vectors,
    )
    return len(chunks)


def retrieve_relevant_uploads(cid, q, top_k=3, max_distance=0.8):
    if upload_collection.count() == 0:
        return []

    query_emb = embeddings.embed_query(q)
    res = upload_collection.query(
        query_embeddings=[query_emb],
        n_results=top_k,
        where={"Cid": cid},
        include=["documents", "distances"],
    )

    docs = res["documents"][0]
    distances = res["distances"][0]

    # only keep chunks that are genuinely close in meaning to the query;
    # discards results that ChromaDB returns just because top_k asked for them
    return [doc for doc, dist in zip(docs, distances) if dist <= max_distance]